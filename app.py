from fastapi import FastAPI, Request
from fastapi.responses import PlainTextResponse, JSONResponse
from fastapi.staticfiles import StaticFiles

import requests
import os
import pandas as pd
from pptx import Presentation
import zipfile
import shutil
import traceback

app = FastAPI()

# ------------------ НАСТРОЙКИ ------------------

BASE_URL = "https://web-production-a9964.up.railway.app"

TEMPLATES_DIR = "templates"
EXCEL_DIR = "excel"
OUTPUT_DIR = "output"
STATE_DIR = "state"
LOG_DIR = "logs"

os.makedirs(TEMPLATES_DIR, exist_ok=True)
os.makedirs(EXCEL_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(STATE_DIR, exist_ok=True)
os.makedirs(LOG_DIR, exist_ok=True)

app.mount("/files", StaticFiles(directory=OUTPUT_DIR), name="files")

# ------------------ ЛОГИРОВАНИЕ ------------------

def save_log(user_id: str, tag: str, data):
    log_path = os.path.join(LOG_DIR, f"{user_id}_{tag}.txt")
    with open(log_path, "w", encoding="utf-8") as f:
        f.write(str(data))

def log_debug(user_id: str, message: str):
    print(f"[{user_id}] {message}")

# ------------------ ИЗВЛЕЧЕНИЕ URL ПО РАСШИРЕНИЮ ------------------

def extract_file_url(variables: list, allowed_extensions: list) -> tuple[str | None, str | None]:
    """
    Ищет URL файла с нужным расширением.
    Возвращает (url, расширение) или (None, None).
    """
    found_urls = []

    for var in variables:
        if not var:
            continue
        payload = var.get("payload") or {}
        url = payload.get("url") or ""
        if not url:
            # попробуем вложенные структуры
            for key, val in payload.items():
                if isinstance(val, str) and val.startswith("http"):
                    url = val
                    break

        if url:
            clean_url = url.split("?")[0].lower()
            for ext in allowed_extensions:
                if clean_url.endswith(ext):
                    found_urls.append((url, ext))

    # берём последний подходящий
    if found_urls:
        return found_urls[-1]

    return None, None

# ------------------ СКАЧИВАНИЕ ФАЙЛА ------------------

def download_file(file_url: str, save_dir: str, force_ext: str = None) -> tuple[str, str]:
    response = requests.get(file_url, timeout=30)
    response.raise_for_status()

    filename = file_url.split("/")[-1].split("?")[0]

    # если имя файла не содержит нужное расширение — принудительно добавляем
    if force_ext and not filename.lower().endswith(force_ext):
        filename = filename + force_ext

    path = os.path.join(save_dir, filename)

    with open(path, "wb") as f:
        f.write(response.content)

    # проверяем что файл реально скачался
    size = os.path.getsize(path)
    print(f"Скачан файл: {path} ({size} байт)")

    return filename, path

# ------------------ ГЕНЕРАЦИЯ PPTX ------------------

def generate_pptx(template_path: str, excel_path: str, user_id: str) -> str:

    # Проверяем шаблон
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Шаблон не найден: {template_path}")

    if not template_path.lower().endswith(".pptx"):
        raise ValueError(f"Шаблон должен быть .pptx, получен: {template_path}")

    # Проверяем Excel
    if not os.path.exists(excel_path):
        raise FileNotFoundError(f"Excel не найден: {excel_path}")

    if not excel_path.lower().endswith(".xlsx"):
        raise ValueError(f"Excel должен быть .xlsx, получен: {excel_path}")

    user_output_dir = os.path.join(OUTPUT_DIR, user_id)

    if os.path.exists(user_output_dir):
        shutil.rmtree(user_output_dir)

    os.makedirs(user_output_dir, exist_ok=True)

    df = pd.read_excel(excel_path)
    print(f"Excel прочитан: {len(df)} строк, колонки: {list(df.columns)}")

    generated_files = []

    for index, row in df.iterrows():
        prs = Presentation(template_path)

        for slide in prs.slides:
            for shape in slide.shapes:

                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            text = run.text
                            for col in df.columns:
                                placeholder = str(col).strip()
                                value = str(row[col])
                                text = text.replace(placeholder, value)
                            run.text = text

        if "Название" in df.columns:
            safe_name = str(row["Название"])
        else:
            safe_name = f"file_{index + 1}"

        bad_chars = ['/', '\\', ':', '*', '?', '"', '<', '>', '|']
        for char in bad_chars:
            safe_name = safe_name.replace(char, "")

        pptx_filename = f"{safe_name}.pptx"
        pptx_path = os.path.join(user_output_dir, pptx_filename)
        prs.save(pptx_path)
        generated_files.append(pptx_path)
        print(f"Создан файл: {pptx_path}")

    zip_name = f"{user_id}_result.zip"
    zip_path = os.path.join(OUTPUT_DIR, zip_name)

    with zipfile.ZipFile(zip_path, "w") as zipf:
        for file in generated_files:
            zipf.write(file, arcname=os.path.basename(file))

    zip_size = os.path.getsize(zip_path)
    print(f"ZIP создан: {zip_path} ({zip_size} байт)")

    return zip_name

# ------------------ ОТЛАДКА ПЕРЕМЕННЫХ ------------------

@app.post("/debug_vars")
async def debug_vars(request: Request):
    """Эндпоинт для отладки — показывает все переменные запроса"""
    try:
        data = await request.json()
        variables = data.get("variables") or []
        contact = data.get("contact") or {}
        user_id = str(contact.get("id", "unknown"))

        save_log(user_id, "debug_vars", data)

        urls_found = []
        for i, var in enumerate(variables):
            if not var:
                continue
            payload = var.get("payload") or {}
            url = payload.get("url")
            if url:
                urls_found.append(f"[{i}] {url}")

        result = f"user_id: {user_id}\n"
        result += f"Найдено URL ({len(urls_found)}):\n"
        result += "\n".join(urls_found) if urls_found else "— нет URL"

        return PlainTextResponse(result)
    except Exception as e:
        return PlainTextResponse(f"Ошибка debug: {str(e)}\n{traceback.format_exc()}")

# ------------------ ЗАГРУЗКА ШАБЛОНА ------------------

@app.post("/upload_template")
async def upload_template(request: Request):
    try:
        data = await request.json()
        variables = data.get("variables") or []
        contact = data.get("contact") or {}
        user_id = str(contact.get("id", "unknown"))

        save_log(user_id, "upload_template_raw", data)
        log_debug(user_id, f"upload_template: переменных {len(variables)}")

        # Ищем ТОЛЬКО .pptx файл
        file_url, ext = extract_file_url(variables, [".pptx"])

        if not file_url:
            # Покажем что вообще пришло для диагностики
            all_urls = []
            for var in variables:
                if not var:
                    continue
                payload = var.get("payload") or {}
                url = payload.get("url")
                if url:
                    all_urls.append(url)

            diag = f"PPTX файл не найден.\n\nВсе URL в запросе:\n"
            diag += "\n".join(all_urls) if all_urls else "— URL не найдены вообще"
            diag += f"\n\nПолный запрос сохранён в лог."
            return PlainTextResponse(diag)

        log_debug(user_id, f"Скачиваем шаблон: {file_url}")
        filename, template_path = download_file(file_url, TEMPLATES_DIR, force_ext=".pptx")

        # Проверяем что это реально PPTX
        try:
            prs = Presentation(template_path)
            slide_count = len(prs.slides)
            log_debug(user_id, f"PPTX валиден, слайдов: {slide_count}")
        except Exception as pptx_err:
            return PlainTextResponse(
                f"Файл скачан, но не является валидным PPTX:\n{str(pptx_err)}\n"
                f"URL: {file_url}"
            )

        state_file = os.path.join(STATE_DIR, f"{user_id}.txt")
        with open(state_file, "w", encoding="utf-8") as f:
            f.write(template_path)

        return PlainTextResponse(
            f"Шаблон загружен ✅\n"
            f"Файл: {filename}\n"
            f"Слайдов: {slide_count}\n\n"
            f"Теперь отправь Excel файл (.xlsx)"
        )

    except Exception as e:
        tb = traceback.format_exc()
        print(f"Ошибка upload_template: {tb}")
        return PlainTextResponse(f"Ошибка upload_template:\n{str(e)}\n\nДетали:\n{tb}")

# ------------------ ЗАГРУЗКА EXCEL ------------------

@app.post("/upload_excel")
async def upload_excel(request: Request):
    try:
        data = await request.json()
        variables = data.get("variables") or []
        contact = data.get("contact") or {}
        user_id = str(contact.get("id", "unknown"))

        save_log(user_id, "upload_excel_raw", data)
        log_debug(user_id, f"upload_excel: переменных {len(variables)}")

        # Проверяем шаблон
        state_file = os.path.join(STATE_DIR, f"{user_id}.txt")
        if not os.path.exists(state_file):
            return PlainTextResponse("Сначала загрузи шаблон PPTX")

        with open(state_file, "r", encoding="utf-8") as f:
            template_path = f.read().strip()

        log_debug(user_id, f"Шаблон из state: {template_path}")

        if not os.path.exists(template_path):
            return PlainTextResponse(
                f"Шаблон не найден на диске: {template_path}\n"
                f"Загрузи шаблон заново."
            )

        # Ищем ТОЛЬКО .xlsx файл
        file_url, ext = extract_file_url(variables, [".xlsx", ".xls"])

        if not file_url:
            all_urls = []
            for var in variables:
                if not var:
                    continue
                payload = var.get("payload") or {}
                url = payload.get("url")
                if url:
                    all_urls.append(url)

            diag = f"Excel файл не найден.\n\nВсе URL в запросе:\n"
            diag += "\n".join(all_urls) if all_urls else "— URL не найдены вообще"
            return PlainTextResponse(diag)

        log_debug(user_id, f"Скачиваем Excel: {file_url}")
        filename, excel_path = download_file(file_url, EXCEL_DIR, force_ext=".xlsx")

        # Проверяем Excel
        try:
            df_test = pd.read_excel(excel_path, nrows=1)
            log_debug(user_id, f"Excel валиден, колонки: {list(df_test.columns)}")
        except Exception as xl_err:
            return PlainTextResponse(
                f"Файл скачан, но не является валидным Excel:\n{str(xl_err)}\n"
                f"URL: {file_url}"
            )

        zip_name = generate_pptx(
            template_path=template_path,
            excel_path=excel_path,
            user_id=user_id
        )

        full_url = f"{BASE_URL}/files/{zip_name}"

        return PlainTextResponse(
            f"Файлы готовы ✅\n\n{full_url}"
        )

    except Exception as e:
        tb = traceback.format_exc()
        print(f"Ошибка upload_excel: {tb}")
        return PlainTextResponse(f"Ошибка сервера:\n{str(e)}\n\nДетали:\n{tb}")

# ------------------ СТАТУС ------------------

@app.get("/", response_class=PlainTextResponse)
def test():
    return "бот работает ✅"

@app.get("/status/{user_id}", response_class=PlainTextResponse)
def status(user_id: str):
    state_file = os.path.join(STATE_DIR, f"{user_id}.txt")
    if os.path.exists(state_file):
        with open(state_file) as f:
            tpl = f.read()
        exists = "✅ найден" if os.path.exists(tpl) else "❌ файл удалён"
        return f"Шаблон: {tpl}\nСтатус: {exists}"
    return "Шаблон не загружен"
