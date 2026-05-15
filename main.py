from fastapi import FastAPI, Header, Request
from fastapi.responses import JSONResponse, PlainTextResponse
from fastapi.staticfiles import StaticFiles

import json
import os
import shutil
import subprocess
import traceback
import zipfile
from pathlib import Path
from typing import Any

import pandas as pd
import requests
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

try:
    import pymorphy3
except Exception:
    pymorphy3 = None

try:
    from pytrovich.detector import PetrovichGenderDetector
    from pytrovich.enums import Case as PetrovichCase
    from pytrovich.enums import Gender as PetrovichGender
    from pytrovich.enums import NamePart
    from pytrovich.maker import PetrovichDeclinationMaker
except Exception:
    PetrovichCase = None
    PetrovichDeclinationMaker = None
    PetrovichGender = None
    PetrovichGenderDetector = None
    NamePart = None


app = FastAPI()

# Public URL of your Railway service. Set this in Railway variables after deploy.
BASE_URL = os.getenv("BASE_URL", "https://web-production-a9964.up.railway.app")

# Token from MAX Master Bot / business.max.ru. Do not paste it into code.
BOT_TOKEN = os.getenv("BOT_TOKEN", "")
MAX_API_URL = os.getenv("MAX_API_URL", "https://platform-api.max.ru")

# Optional secret for MAX webhook. Set the same value in Railway and webhook subscription.
WEBHOOK_SECRET = os.getenv("WEBHOOK_SECRET", "")
STRICT_WEBHOOK_SECRET = os.getenv("STRICT_WEBHOOK_SECRET", "false").lower() == "true"
INCLUDE_DEBUG_PPTX = os.getenv("INCLUDE_DEBUG_PPTX", "false").lower() == "true"

TEMPLATES_DIR = Path("templates")
EXCEL_DIR = Path("excel")
OUTPUT_DIR = Path("output")
STATE_DIR = Path("state")

for directory in (TEMPLATES_DIR, EXCEL_DIR, OUTPUT_DIR, STATE_DIR):
    directory.mkdir(exist_ok=True)

app.mount("/files", StaticFiles(directory=str(OUTPUT_DIR)), name="files")

MORPH = None
PETROVICH_MAKER = None
PETROVICH_DETECTOR = None

CASES = {
    "nomn": {"name": "Именительный", "question": "кто? что?"},
    "gent": {"name": "Родительный", "question": "кого? чего?"},
    "datv": {"name": "Дательный", "question": "кому? чему?"},
    "accs": {"name": "Винительный", "question": "кого? что?"},
    "ablt": {"name": "Творительный", "question": "кем? чем?"},
    "loct": {"name": "Предложный", "question": "о ком? о чем?"},
}

PETROVICH_CASES = {
    "gent": "GENITIVE",
    "datv": "DATIVE",
    "accs": "ACCUSATIVE",
    "ablt": "INSTRUMENTAL",
    "loct": "PREPOSITIONAL",
}


# ------------------ BASIC HELPERS ------------------

def sanitize_filename(value: Any, fallback: str = "file") -> str:
    name = str(value or fallback).strip()
    bad_chars = ['/', '\\', ':', '*', '?', '"', '<', '>', '|']
    for char in bad_chars:
        name = name.replace(char, "")
    return name[:120] or fallback


def user_templates_dir(user_id: str) -> Path:
    path = TEMPLATES_DIR / sanitize_filename(user_id, fallback="default")
    path.mkdir(exist_ok=True)
    return path


def state_path(user_id: str) -> Path:
    return STATE_DIR / f"{sanitize_filename(user_id, fallback='default')}.json"


def normalize_state(state: dict) -> dict:
    state.setdefault("templates", [])

    # Migration from the first one-template version.
    if state.get("template_path") and not state["templates"]:
        state["templates"].append({
            "name": state.get("template_name") or Path(state["template_path"]).name,
            "path": state["template_path"],
        })
        state.pop("template_path", None)
        state.pop("template_name", None)

    return state


def load_state(user_id: str) -> dict:
    path = state_path(user_id)
    if not path.exists():
        return normalize_state({})
    try:
        return normalize_state(json.loads(path.read_text(encoding="utf-8")))
    except Exception:
        return normalize_state({})


def save_state(user_id: str, data: dict) -> None:
    state_path(user_id).write_text(
        json.dumps(normalize_state(data), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def unique_path(directory: Path, filename: str) -> Path:
    filename = sanitize_filename(filename)
    path = directory / filename
    if not path.exists():
        return path

    stem = path.stem
    suffix = path.suffix
    counter = 2
    while True:
        candidate = directory / f"{stem}_{counter}{suffix}"
        if not candidate.exists():
            return candidate
        counter += 1


def download_file(file_url: str, save_dir: Path, filename: str | None = None, force_ext: str | None = None):
    response = requests.get(file_url, timeout=60)
    response.raise_for_status()

    if not filename:
        filename = file_url.split("/")[-1].split("?")[0] or "file"

    filename = sanitize_filename(filename)
    if force_ext and not filename.lower().endswith(force_ext):
        filename = f"{filename}{force_ext}"

    path = unique_path(save_dir, filename)
    path.write_bytes(response.content)
    return path.name, path


def format_value(value: Any) -> str:
    if pd.isna(value):
        return ""
    return str(value)


def template_list_text(templates: list[dict]) -> str:
    if not templates:
        return "У тебя пока нет загруженных шаблонов."

    lines = ["Твои шаблоны:"]
    for index, template in enumerate(templates, start=1):
        lines.append(f"{index}. {template.get('name') or 'template.pptx'}")
    return "\n".join(lines)


def delete_template_file(template: dict) -> None:
    path = Path(template.get("path") or "")
    try:
        if path.exists() and path.is_file():
            path.unlink()
    except Exception:
        pass


# ------------------ RUSSIAN CASES ------------------

def get_morph():
    global MORPH
    if MORPH is not None:
        return MORPH
    if pymorphy3 is None:
        raise RuntimeError("Библиотека pymorphy3 не установлена. Проверь requirements.txt и redeploy.")
    MORPH = pymorphy3.MorphAnalyzer()
    return MORPH


def get_petrovich():
    global PETROVICH_MAKER, PETROVICH_DETECTOR
    if PetrovichDeclinationMaker is None or PetrovichGenderDetector is None:
        return None, None
    if PETROVICH_MAKER is None:
        PETROVICH_MAKER = PetrovichDeclinationMaker()
    if PETROVICH_DETECTOR is None:
        PETROVICH_DETECTOR = PetrovichGenderDetector()
    return PETROVICH_MAKER, PETROVICH_DETECTOR


def preserve_case(original: str, changed: str) -> str:
    if original.isupper():
        return changed.upper()
    if original[:1].isupper():
        return changed.capitalize()
    return changed


def infer_gender_from_fio_parts(parts: list[str]) -> str | None:
    for part in reversed(parts):
        clean = part.strip("- ").lower()
        if clean.endswith(("вна", "чна", "кызы")):
            return "femn"
        if clean.endswith(("вич", "ич", "оглы")):
            return "masc"

    morph = get_morph()
    votes = {"femn": 0, "masc": 0}
    for part in parts:
        clean = part.strip("- ")
        if not clean:
            continue
        parsed = morph.parse(clean)[0]
        if "femn" in parsed.tag:
            votes["femn"] += 1
        if "masc" in parsed.tag:
            votes["masc"] += 1

    if votes["femn"] > votes["masc"]:
        return "femn"
    if votes["masc"] > votes["femn"]:
        return "masc"
    return None


def petrovich_gender(gender: str | None):
    if PetrovichGender is None:
        return None
    if gender == "femn":
        return PetrovichGender.FEMALE
    if gender == "masc":
        return PetrovichGender.MALE
    return None


def petrovich_case(case_code: str):
    if PetrovichCase is None:
        return None
    case_name = PETROVICH_CASES.get(case_code)
    if not case_name:
        return None
    return getattr(PetrovichCase, case_name, None)


def detect_petrovich_gender(parts: list[str], gender: str | None):
    maker, detector = get_petrovich()
    explicit_gender = petrovich_gender(gender)
    if explicit_gender or detector is None:
        return explicit_gender

    lastname = parts[0] if len(parts) >= 1 else None
    firstname = parts[1] if len(parts) >= 2 else None
    middlename = " ".join(parts[2:]) if len(parts) >= 3 else None

    try:
        detected = detector.detect(
            lastname=lastname,
            firstname=firstname,
            middlename=middlename,
        )
        if detected:
            return detected
    except Exception:
        return None

    return None


def inflect_fio_with_petrovich(parts: list[str], case_code: str, gender: str | None) -> str | None:
    maker, _ = get_petrovich()
    target_case = petrovich_case(case_code)
    target_gender = detect_petrovich_gender(parts, gender)

    if maker is None or target_case is None or target_gender is None:
        return None

    name_parts = []
    if len(parts) >= 1:
        name_parts.append((NamePart.LASTNAME, parts[0]))
    if len(parts) >= 2:
        name_parts.append((NamePart.FIRSTNAME, parts[1]))
    if len(parts) >= 3:
        name_parts.append((NamePart.MIDDLENAME, " ".join(parts[2:])))

    try:
        changed = [
            preserve_case(original, maker.make(name_part, target_gender, target_case, original))
            for name_part, original in name_parts
        ]
        return " ".join(changed)
    except Exception:
        return None


def inflect_word(word: str, case_code: str, gender: str | None = None) -> str:
    if case_code == "nomn" or not word:
        return word

    morph = get_morph()
    parses = morph.parse(word)
    parsed = parses[0]

    if gender:
        gender_parses = [item for item in parses if gender in item.tag]
        if gender_parses:
            parsed = gender_parses[0]

    required = {case_code}
    if gender:
        required.add(gender)

    inflected = parsed.inflect(required)
    if not inflected and gender:
        inflected = parsed.inflect({case_code})

    if not inflected:
        return word
    return preserve_case(word, inflected.word)


def inflect_name_part(part: str, case_code: str, gender: str | None = None) -> str:
    if "-" in part:
        return "-".join(inflect_word(piece, case_code, gender) for piece in part.split("-"))
    return inflect_word(part, case_code, gender)


def inflect_fio(value: Any, case_code: str) -> str:
    text = format_value(value).strip()
    if case_code == "nomn" or not text:
        return text

    parts = text.split()
    gender = infer_gender_from_fio_parts(parts)
    petrovich_result = inflect_fio_with_petrovich(parts, case_code, gender)
    if petrovich_result:
        return petrovich_result

    return " ".join(inflect_name_part(part, case_code, gender) for part in parts)


def normalize_placeholder_name(value: Any) -> str:
    name = str(value).strip().lower()
    name = name.strip("%{}[]() ")
    return name.replace(".", "").replace(" ", "")


def should_inflect_column(column: Any) -> bool:
    name = normalize_placeholder_name(column)
    return name in {
        "фио",
        "фамилияимяотчество",
        "фамилия",
        "имя",
        "отчество",
        "полноеимя",
    }


def row_value_for_column(row: pd.Series, column: Any, case_code: str) -> str:
    value = row[column]
    if should_inflect_column(column):
        return inflect_fio(value, case_code)
    return format_value(value)


# ------------------ POWERPOINT AND PDF GENERATION ------------------

def placeholder_for_column(column: Any) -> str:
    name = str(column).strip()
    if name.startswith("%") and name.endswith("%"):
        return name
    return f"%{name}%"


def build_replacements(row: pd.Series, columns, case_code: str) -> list[tuple[str, str]]:
    replacements = []
    for column in columns:
        placeholder = placeholder_for_column(column)
        replacements.append((placeholder, row_value_for_column(row, column, case_code)))

    # Longer placeholders first prevents partial replacement surprises.
    replacements.sort(key=lambda item: len(item[0]), reverse=True)
    return replacements


def is_single_placeholder_text(text: str, replacements: list[tuple[str, str]]) -> bool:
    clean_text = text.strip()
    if not clean_text:
        return False
    return any(clean_text == placeholder for placeholder, _ in replacements)


def replace_text_in_shape(shape, row: pd.Series, columns, case_code: str) -> None:
    if not shape.has_text_frame:
        return

    replacements = build_replacements(row, columns, case_code)

    for paragraph in shape.text_frame.paragraphs:
        original_full_text = "".join(run.text for run in paragraph.runs)
        if is_single_placeholder_text(original_full_text, replacements) and paragraph.alignment is None:
            paragraph.alignment = PP_ALIGN.CENTER

        # First try replacing inside each run. This preserves PowerPoint formatting
        # much better when the placeholder is not split across multiple runs.
        replaced_inside_runs = False
        for run in paragraph.runs:
            original_text = run.text
            new_text = original_text
            for placeholder, value in replacements:
                if placeholder in new_text:
                    new_text = new_text.replace(placeholder, value)

            if new_text != original_text:
                run.text = new_text
                replaced_inside_runs = True

        if replaced_inside_runs:
            continue

        # Fallback for rare cases where PowerPoint split a placeholder across runs.
        full_text = "".join(run.text for run in paragraph.runs)
        if not full_text:
            continue

        replaced = False
        for placeholder, value in replacements:
            if placeholder in full_text:
                full_text = full_text.replace(placeholder, value)
                replaced = True

        if replaced and paragraph.runs:
            for run in paragraph.runs:
                run.text = ""
            paragraph.runs[0].text = full_text


def libreoffice_binary() -> str:
    binary = shutil.which("libreoffice") or shutil.which("soffice")
    if not binary:
        raise RuntimeError(
            "Не найден LibreOffice для конвертации PPTX в PDF. "
            "В проект добавлен nixpacks.toml, после загрузки на Railway сделай redeploy."
        )
    return binary


def convert_pptx_to_pdf(pptx_path: Path, output_dir: Path) -> Path:
    lo_profile = Path("/tmp/lo_profile_fixed")
    lo_fonts_dir = lo_profile / "user" / "fonts"
    lo_fonts_dir.mkdir(parents=True, exist_ok=True)

    # Копируем Circe прямо в профиль LibreOffice — он точно их увидит
    app_fonts = Path("/app/fonts")
    if app_fonts.exists():
        for font_file in app_fonts.iterdir():
            if font_file.suffix.lower() in (".ttf", ".otf"):
                dest = lo_fonts_dir / font_file.name
                if not dest.exists():
                    shutil.copy2(font_file, dest)

    command = [
        libreoffice_binary(),
        "--headless",
        "--norestore",
        "--nofirststartwizard",
        f"-env:UserInstallation=file://{lo_profile}",
        "--convert-to", "pdf:impress_pdf_Export",
        "--outdir", str(output_dir),
        str(pptx_path),
    ]
    env = os.environ.copy()
    env["HOME"] = "/tmp"
    env["TMPDIR"] = "/tmp"

    result = subprocess.run(command, capture_output=True, text=True, timeout=180, env=env)
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice error:\n{result.stderr or result.stdout}")

    pdf_path = output_dir / f"{pptx_path.stem}.pdf"
    if not pdf_path.exists():
        raise RuntimeError("PDF не появился после конвертации.")
    return pdf_path

from pptx.enum.text import MSO_AUTO_SIZE

KNOWN_SYSTEM_FONTS = {
    "liberation sans", "liberation serif", "liberation mono",
    "dejavu sans", "dejavu serif", "noto sans", "noto serif",
    "lato", "open sans", "roboto", "montserrat", "pt sans", "pt serif",
    "arial", "times new roman", "courier new", "calibri",
}
FALLBACK_FONT = "Liberation Sans"

def substitute_unknown_fonts(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            font_substituted = False

            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    font_name = (run.font.name or "").strip()
                    if font_name and font_name.lower() not in KNOWN_SYSTEM_FONTS:
                        run.font.name = FALLBACK_FONT
                        font_substituted = True

            if font_substituted:
                # Включаем авто-уменьшение шрифта, чтобы текст не вылезал за пределы блока
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                tf.word_wrap = True
                            
def generate_pptx_with_replacements(template_path: Path, row: pd.Series, columns, case_code: str, output_path: Path) -> None:
    """Заменяет текст через python-pptx, но пишет только слайды в оригинальный ZIP.
    Вся остальная структура (темы, шрифты, лэйауты) берётся из шаблона как есть."""
    
    # Используем python-pptx только для замены текста
    prs = Presentation(str(template_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, row, columns, case_code)
    
    # Получаем изменённые XML только слайдов
    modified_slides = {}
    for slide in prs.slides:
        zip_name = slide.part.partname.lstrip('/')   # e.g. ppt/slides/slide1.xml
        modified_slides[zip_name] = slide.part.blob  # XML байты только этого слайда
    
    # Читаем оригинальный шаблон как ZIP
    with zipfile.ZipFile(template_path, 'r') as zin:
        file_contents = {name: zin.read(name) for name in zin.namelist()}
    
    # Патчим только слайды, всё остальное оставляем оригинальным
    for zip_name, xml_bytes in modified_slides.items():
        if zip_name in file_contents:
            file_contents[zip_name] = xml_bytes
    
    # Пишем новый ZIP
    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name, data in file_contents.items():
            zout.writestr(name, data)


def generate_pdf_zip(template_path: Path, excel_path: Path, user_id: str, case_code: str) -> str:
    user_output_dir = OUTPUT_DIR / sanitize_filename(user_id, fallback="default")

    if user_output_dir.exists():
        shutil.rmtree(user_output_dir)
    user_output_dir.mkdir(exist_ok=True)

    df = pd.read_excel(excel_path)
    if df.empty:
        raise ValueError("Excel пустой. Добавь хотя бы одну строку с данными.")

    generated_pdfs = []

    for index, row in df.iterrows():
        safe_name = sanitize_filename(row[df.columns[0]], fallback=f"presentation_{index + 1}")
        pptx_path = unique_path(user_output_dir, f"{safe_name}.pptx")
        
        # Патчим ZIP напрямую вместо prs.save()
        generate_pptx_with_replacements(template_path, row, df.columns, case_code, pptx_path)
        
        pdf_path = convert_pptx_to_pdf(pptx_path, user_output_dir)
        generated_pdfs.append(pdf_path)

    zip_name = f"{sanitize_filename(user_id, fallback='default')}_result.zip"
    zip_path = OUTPUT_DIR / zip_name

    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zipf:
        for file in generated_pdfs:
            zipf.write(file, arcname=file.name)

    return zip_name
# ------------------ MAX API HELPERS ------------------

def max_headers() -> dict:
    if not BOT_TOKEN:
        raise RuntimeError("BOT_TOKEN не задан в Railway Variables")
    return {"Authorization": BOT_TOKEN, "Content-Type": "application/json"}


def inline_keyboard(buttons: list[list[dict]]) -> dict:
    return {"type": "inline_keyboard", "payload": {"buttons": buttons}}


def main_menu_keyboard() -> dict:
    return inline_keyboard([
        [{"type": "callback", "text": "Мои шаблоны", "payload": "my_templates"}],
        [{"type": "callback", "text": "Сгенерировать", "payload": "generate"}],
        [{"type": "callback", "text": "Инструкция", "payload": "instruction"}],
    ])


def templates_keyboard() -> dict:
    return inline_keyboard([
        [{"type": "callback", "text": "Добавить шаблон", "payload": "add_template"}],
        [{"type": "callback", "text": "Удалить шаблон", "payload": "delete_template"}],
        [{"type": "callback", "text": "Главное меню", "payload": "main_menu"}],
    ])


def back_to_menu_keyboard() -> dict:
    return inline_keyboard([
        [{"type": "callback", "text": "Главное меню", "payload": "main_menu"}],
    ])


def case_keyboard() -> dict:
    rows = []
    for code, info in CASES.items():
        rows.append([{"type": "callback", "text": info["name"], "payload": f"case_{code}"}])
    rows.append([{"type": "callback", "text": "Главное меню", "payload": "main_menu"}])
    return inline_keyboard(rows)


def send_max_message(user_id: str, text: str, attachments: list[dict] | None = None) -> None:
    payload = {"text": text, "attachments": attachments or [], "link": None}
    response = requests.post(
        f"{MAX_API_URL}/messages",
        params={"user_id": user_id},
        headers=max_headers(),
        json=payload,
        timeout=30,
    )
    if response.ok:
        return

    print(f"MAX send message error {response.status_code}: {response.text}")

    if attachments:
        fallback_response = requests.post(
            f"{MAX_API_URL}/messages",
            params={"user_id": user_id},
            headers=max_headers(),
            json={"text": text, "attachments": [], "link": None},
            timeout=30,
        )
        fallback_response.raise_for_status()
        return

    response.raise_for_status()


def answer_callback(callback_id: str, notification: str = "Готово") -> None:
    if not callback_id:
        return

    response = requests.post(
        f"{MAX_API_URL}/answers",
        params={"callback_id": callback_id},
        headers=max_headers(),
        json={"notification": notification, "message": None},
        timeout=30,
    )
    response.raise_for_status()


def extract_user_id(update: dict) -> str | None:
    if update.get("update_type") == "bot_started":
        return str((update.get("user") or {}).get("user_id") or "")

    if update.get("update_type") == "message_callback":
        return str(((update.get("callback") or {}).get("user") or {}).get("user_id") or "")

    message = update.get("message") or {}
    sender = message.get("sender") or {}
    return str(sender.get("user_id") or "")


def extract_text(update: dict) -> str:
    message = update.get("message") or {}
    body = message.get("body") or {}
    return (body.get("text") or "").strip()


def extract_file_attachment(update: dict, allowed_extensions: list[str]) -> tuple[str | None, str | None, str | None]:
    message = update.get("message") or {}
    body = message.get("body") or {}
    attachments = body.get("attachments") or []

    found = []
    for attachment in attachments:
        if not isinstance(attachment, dict) or attachment.get("type") != "file":
            continue

        payload = attachment.get("payload") or {}
        url = payload.get("url") or attachment.get("url") or ""
        filename = attachment.get("filename") or url.split("/")[-1].split("?")[0]
        clean_name = filename.lower()
        clean_url = url.split("?")[0].lower()

        for ext in allowed_extensions:
            if clean_name.endswith(ext) or clean_url.endswith(ext):
                found.append((url, filename, ext))

    return found[-1] if found else (None, None, None)


# ------------------ MAX BOT FLOW ------------------

def send_main_menu(user_id: str) -> None:
    send_max_message(user_id, "Выбери действие:", [main_menu_keyboard()])


def send_instruction(user_id: str) -> None:
    send_max_message(
        user_id,
        "Инструкция:\n\n"
        "1. Сначала открой «Мои шаблоны» и добавь PPTX-шаблон.\n"
        "2. В шаблоне напиши плейсхолдеры через проценты: %ФИО%, %КЛАСС%, %ДАТА%, %НОМИНАЦИЯ%.\n"
        "3. В Excel сделай колонки с такими же названиями: %ФИО%, %КЛАСС%, %ДАТА%, %НОМИНАЦИЯ%.\n"
        "4. Каждая строка Excel станет отдельной грамотой.\n"
        "5. Нажми «Сгенерировать», выбери номер шаблона, отправь Excel и выбери падеж.\n"
        "6. В конце бот пришлет ссылку на ZIP-архив с PDF-файлами.\n\n"
        "Для склонения ФИО используй колонку %ФИО%. "
        "Автоматические падежи могут ошибаться на редких или сложных ФИО.",
        [back_to_menu_keyboard()],
    )


def selected_template_from_state(state: dict) -> dict | None:
    templates = state.get("templates") or []
    index = state.get("selected_template_index")
    if not isinstance(index, int) or index < 0 or index >= len(templates):
        return None
    return templates[index]


async def handle_callback(update: dict, user_id: str) -> None:
    callback = update.get("callback") or {}
    payload = callback.get("payload")
    callback_id = callback.get("callback_id")
    state = load_state(user_id)

    if payload == "main_menu":
        state["mode"] = None
        save_state(user_id, state)
        answer_callback(callback_id, "Меню")
        send_main_menu(user_id)
        return

    if payload == "my_templates":
        state["mode"] = None
        save_state(user_id, state)
        answer_callback(callback_id, "Шаблоны")
        send_max_message(user_id, template_list_text(state["templates"]), [templates_keyboard()])
        return

    if payload == "add_template":
        state["mode"] = "waiting_template"
        save_state(user_id, state)
        answer_callback(callback_id, "Жду PPTX")
        send_max_message(user_id, "Отправь PPTX-шаблон одним файлом.", [back_to_menu_keyboard()])
        return

    if payload == "delete_template":
        if not state["templates"]:
            answer_callback(callback_id, "Нет шаблонов")
            send_max_message(user_id, "Удалять пока нечего: шаблоны еще не загружены.", [back_to_menu_keyboard()])
            return

        state["mode"] = "waiting_delete_template_number"
        save_state(user_id, state)
        answer_callback(callback_id, "Жду номер")
        send_max_message(
            user_id,
            f"{template_list_text(state['templates'])}\n\n"
            "Отправь цифру шаблона, который нужно удалить.",
            [back_to_menu_keyboard()],
        )
        return

    if payload == "generate":
        if not state["templates"]:
            answer_callback(callback_id, "Нет шаблонов")
            send_max_message(
                user_id,
                "Сначала добавь хотя бы один PPTX-шаблон в разделе «Мои шаблоны».",
                [templates_keyboard()],
            )
            return

        state["mode"] = "waiting_generate_template_number"
        save_state(user_id, state)
        answer_callback(callback_id, "Выбери шаблон")
        send_max_message(
            user_id,
            f"{template_list_text(state['templates'])}\n\n"
            "Отправь цифру нужного шаблона.",
            [back_to_menu_keyboard()],
        )
        return

    if payload == "instruction":
        answer_callback(callback_id, "Инструкция")
        send_instruction(user_id)
        return

    if payload and payload.startswith("case_"):
        case_code = payload.replace("case_", "", 1)
        if case_code not in CASES:
            answer_callback(callback_id, "Неизвестный падеж")
            return

        template = selected_template_from_state(state)
        excel_path = state.get("excel_path")
        if not template or not excel_path:
            answer_callback(callback_id, "Нет файлов")
            send_max_message(user_id, "Не хватает шаблона или Excel. Начни генерацию заново.", [main_menu_keyboard()])
            return

        template_path = Path(template["path"])
        if not template_path.exists() or not Path(excel_path).exists():
            answer_callback(callback_id, "Файл не найден")
            send_max_message(user_id, "Один из файлов не найден. Загрузи шаблон и Excel заново.", [main_menu_keyboard()])
            return

        answer_callback(callback_id, "Генерирую")
        send_max_message(
            user_id,
            f"Падеж выбран: {CASES[case_code]['name']}.\n"
            "Генерирую PDF-файлы и собираю ZIP. Это может занять немного времени.",
        )

        zip_name = generate_pdf_zip(template_path, Path(excel_path), user_id, case_code)
        state["mode"] = None
        state.pop("excel_path", None)
        state.pop("selected_template_index", None)
        save_state(user_id, state)

        full_url = f"{BASE_URL}/files/{zip_name}"
        send_max_message(user_id, f"Файлы готовы.\n\nСкачать ZIP:\n{full_url}", [main_menu_keyboard()])


async def handle_message(update: dict, user_id: str) -> None:
    text = extract_text(update)
    text_lower = text.lower()

    if text_lower in {"/start", "старт", "меню", "главное меню"}:
        state = load_state(user_id)
        state["mode"] = None
        save_state(user_id, state)
        send_main_menu(user_id)
        return

    state = load_state(user_id)
    mode = state.get("mode")

    if mode == "waiting_template":
        file_url, filename, _ = extract_file_attachment(update, [".pptx"])
        if not file_url:
            send_max_message(user_id, "Я не вижу PPTX-файл. Отправь именно файл с расширением .pptx.", [back_to_menu_keyboard()])
            return

        _, template_path = download_file(
            file_url,
            user_templates_dir(user_id),
            filename=filename,
            force_ext=".pptx",
        )

        try:
            prs = Presentation(str(template_path))
            slide_count = len(prs.slides)
        except Exception as exc:
            send_max_message(user_id, f"Файл не похож на нормальный PPTX:\n{exc}", [back_to_menu_keyboard()])
            return

        state["templates"].append({"name": template_path.name, "path": str(template_path)})
        state["mode"] = None
        save_state(user_id, state)

        send_max_message(
            user_id,
            f"Шаблон добавлен.\nФайл: {template_path.name}\nСлайдов: {slide_count}",
            [back_to_menu_keyboard()],
        )
        return

    if mode == "waiting_delete_template_number":
        if not text.isdigit():
            send_max_message(user_id, "Отправь только цифру шаблона из списка.", [back_to_menu_keyboard()])
            return

        index = int(text) - 1
        templates = state["templates"]
        if index < 0 or index >= len(templates):
            send_max_message(user_id, "Такого номера нет. Проверь список и отправь цифру ещё раз.", [back_to_menu_keyboard()])
            return

        removed = templates.pop(index)
        delete_template_file(removed)
        state["mode"] = None
        save_state(user_id, state)

        send_max_message(user_id, f"Шаблон удалён: {removed.get('name')}", [back_to_menu_keyboard()])
        return

    if mode == "waiting_generate_template_number":
        if not text.isdigit():
            send_max_message(user_id, "Отправь только цифру нужного шаблона.", [back_to_menu_keyboard()])
            return

        index = int(text) - 1
        templates = state["templates"]
        if index < 0 or index >= len(templates):
            send_max_message(user_id, "Такого номера нет. Проверь список и отправь цифру ещё раз.", [back_to_menu_keyboard()])
            return

        template_path = Path(templates[index]["path"])
        if not template_path.exists():
            send_max_message(user_id, "Этот шаблон не найден на сервере. Удали его и загрузи заново.", [templates_keyboard()])
            return

        state["selected_template_index"] = index
        state["mode"] = "waiting_excel"
        save_state(user_id, state)

        send_max_message(
            user_id,
            f"Выбран шаблон: {templates[index]['name']}.\n\n"
            "Теперь отправь Excel-файл .xlsx или .xls.",
            [back_to_menu_keyboard()],
        )
        return

    if mode == "waiting_excel":
        file_url, filename, ext = extract_file_attachment(update, [".xlsx", ".xls"])
        if not file_url:
            send_max_message(user_id, "Я не вижу Excel-файл. Отправь .xlsx или .xls.", [back_to_menu_keyboard()])
            return

        user_excel_dir = EXCEL_DIR / sanitize_filename(user_id, fallback="default")
        user_excel_dir.mkdir(exist_ok=True)
        _, excel_path = download_file(file_url, user_excel_dir, filename=filename, force_ext=ext or ".xlsx")

        state["excel_path"] = str(excel_path)
        state["mode"] = "waiting_case"
        save_state(user_id, state)

        send_max_message(user_id, "Файл Excel загружен, выбери нужный падеж:", [case_keyboard()])
        return

    if mode == "waiting_case":
        send_max_message(user_id, "Выбери падеж кнопкой под предыдущим сообщением.", [case_keyboard()])
        return

    send_main_menu(user_id)


async def handle_max_update(update: dict) -> None:
    user_id = extract_user_id(update)
    if not user_id:
        return

    update_type = update.get("update_type")

    if update_type == "bot_started":
        send_main_menu(user_id)
        return

    if update_type == "message_callback":
        await handle_callback(update, user_id)
        return

    if update_type == "message_created":
        await handle_message(update, user_id)


# ------------------ MAX WEBHOOK ENDPOINTS ------------------

@app.post("/max/webhook")
async def max_webhook(request: Request, x_max_bot_api_secret: str | None = Header(default=None)):
    try:
        if WEBHOOK_SECRET and x_max_bot_api_secret and x_max_bot_api_secret != WEBHOOK_SECRET:
            return JSONResponse({"ok": False, "error": "bad secret"}, status_code=403)

        if WEBHOOK_SECRET and not x_max_bot_api_secret:
            print("MAX webhook warning: secret header is missing")
            if STRICT_WEBHOOK_SECRET:
                return JSONResponse({"ok": False, "error": "missing secret"}, status_code=403)

        update = await request.json()
        print(f"MAX webhook update_type: {update.get('update_type')}")
        await handle_max_update(update)
        return {"ok": True}
    except Exception as exc:
        print("MAX webhook error:", exc)
        print(traceback.format_exc())
        return {"ok": True}


def register_max_webhook() -> str:
    webhook_url = f"{BASE_URL}/max/webhook"
    body = {
        "url": webhook_url,
        "update_types": ["bot_started", "message_created", "message_callback"],
    }
    if WEBHOOK_SECRET:
        body["secret"] = WEBHOOK_SECRET

    response = requests.post(
        f"{MAX_API_URL}/subscriptions",
        headers=max_headers(),
        json=body,
        timeout=30,
    )
    return f"MAX response {response.status_code}:\n{response.text}"


@app.get("/setup_max_webhook", response_class=PlainTextResponse)
def setup_max_webhook_from_browser():
    return register_max_webhook()


@app.post("/setup_max_webhook", response_class=PlainTextResponse)
def setup_max_webhook():
    return register_max_webhook()


# ------------------ OLD WATBOT COMPATIBILITY ENDPOINTS ------------------

def extract_watbot_file_url(variables: list, allowed_extensions: list[str]):
    found_urls = []

    for var in variables:
        if not var:
            continue

        if isinstance(var, dict):
            payload = var.get("payload") or {}
            url = payload.get("url") or var.get("url") or ""
        elif isinstance(var, str):
            url = var
        else:
            url = ""

        if url:
            clean_url = url.split("?")[0].lower()
            for ext in allowed_extensions:
                if clean_url.endswith(ext):
                    found_urls.append((url, ext))

    return found_urls[-1] if found_urls else (None, None)


@app.post("/upload_template")
async def upload_template(request: Request):
    try:
        data = await request.json()
        variables = data if isinstance(data, list) else data.get("variables") or []
        contact = {} if isinstance(data, list) else data.get("contact") or {}
        user_id = str(contact.get("id", "default"))

        file_url, _ = extract_watbot_file_url(variables, [".pptx"])
        if not file_url:
            return PlainTextResponse("PPTX не найден.")

        filename, template_path = download_file(file_url, user_templates_dir(user_id), force_ext=".pptx")

        try:
            prs = Presentation(str(template_path))
            slide_count = len(prs.slides)
        except Exception as exc:
            return PlainTextResponse(f"Файл не является валидным PPTX:\n{exc}")

        state = load_state(user_id)
        state["templates"].append({"name": filename, "path": str(template_path)})
        state["mode"] = None
        save_state(user_id, state)

        return PlainTextResponse(
            f"Шаблон загружен.\n"
            f"Файл: {filename} ({slide_count} слайдов)\n\n"
            f"Теперь отправь Excel файл (.xlsx)"
        )

    except Exception as exc:
        return PlainTextResponse(f"Ошибка upload_template:\n{exc}\n\n{traceback.format_exc()}")


@app.post("/upload_excel")
async def upload_excel(request: Request):
    try:
        data = await request.json()
        variables = data if isinstance(data, list) else data.get("variables") or []
        contact = {} if isinstance(data, list) else data.get("contact") or {}
        user_id = str(contact.get("id", "default"))

        state = load_state(user_id)
        if not state["templates"]:
            return PlainTextResponse("Сначала загрузи шаблон PPTX")

        template_path = Path(state["templates"][-1]["path"])
        if not template_path.exists():
            return PlainTextResponse("Шаблон не найден на диске. Загрузи шаблон заново.")

        file_url, ext = extract_watbot_file_url(variables, [".xlsx", ".xls"])
        if not file_url:
            return PlainTextResponse("Excel не найден.")

        user_excel_dir = EXCEL_DIR / sanitize_filename(user_id, fallback="default")
        user_excel_dir.mkdir(exist_ok=True)
        _, excel_path = download_file(file_url, user_excel_dir, force_ext=ext or ".xlsx")
        zip_name = generate_pdf_zip(template_path, excel_path, user_id, "nomn")

        full_url = f"{BASE_URL}/files/{zip_name}"
        return PlainTextResponse(f"Файлы готовы.\n\n{full_url}")

    except Exception as exc:
        return PlainTextResponse(f"Ошибка сервера:\n{exc}\n\n{traceback.format_exc()}")


# ------------------ STATUS ------------------

@app.get("/", response_class=PlainTextResponse)
def status():
    return "бот работает"


@app.get("/debug", response_class=PlainTextResponse)
def debug():
    libreoffice = shutil.which("libreoffice") or shutil.which("soffice") or "не найден"
    return (
        "debug\n"
        f"BOT_TOKEN: {'задан' if BOT_TOKEN else 'не задан'}\n"
        f"BASE_URL: {BASE_URL}\n"
        f"WEBHOOK_SECRET: {'задан' if WEBHOOK_SECRET else 'не задан'}\n"
        f"STRICT_WEBHOOK_SECRET: {STRICT_WEBHOOK_SECRET}\n"
        f"MAX_API_URL: {MAX_API_URL}\n"
        f"pymorphy3: {'установлен' if pymorphy3 else 'не установлен'}\n"
        f"pytrovich: {'установлен' if PetrovichDeclinationMaker else 'не установлен'}\n"
        f"INCLUDE_DEBUG_PPTX: {INCLUDE_DEBUG_PPTX}\n"
        f"LibreOffice: {libreoffice}\n"
    )


@app.get("/debug_inflect", response_class=PlainTextResponse)
def debug_inflect(fio: str, case: str = "datv"):
    if case not in CASES:
        return f"Неизвестный падеж: {case}"
    return inflect_fio(fio, case)

@app.get("/debug_fonts", response_class=PlainTextResponse)
def debug_fonts():
    import subprocess
    lines = []
    
    # Что видит система
    fc = subprocess.run(["fc-list"], capture_output=True, text=True)
    circe_lines = [l for l in fc.stdout.splitlines() if "irce" in l]
    lines.append(f"=== Circe в fc-list ({len(circe_lines)} шт) ===")
    lines.extend(circe_lines or ["НЕ НАЙДЕН"])
    
    # Что есть в папке /app/fonts
    lines.append("\n=== /app/fonts ===")
    p = Path("/app/fonts")
    if p.exists():
        lines.extend(str(f) for f in p.iterdir())
    else:
        lines.append("ПАПКА НЕ СУЩЕСТВУЕТ")
    
    return "\n".join(lines)

@app.get("/debug_lo_fonts", response_class=PlainTextResponse)
def debug_lo_fonts():
    import subprocess
    lines = []
    
    # Где реально лежит LibreOffice
    lo = shutil.which("libreoffice") or shutil.which("soffice") or "не найден"
    lines.append(f"LibreOffice: {lo}")
    
    # Папка шрифтов внутри LibreOffice
    lo_font_dirs = [
        "/usr/lib/libreoffice/share/fonts/truetype",
        "/usr/lib64/libreoffice/share/fonts/truetype",
        "/opt/libreoffice/share/fonts/truetype",
    ]
    for d in lo_font_dirs:
        p = Path(d)
        if p.exists():
            files = list(p.rglob("*.ttf"))
            circe = [f for f in files if "irce" in f.name]
            lines.append(f"\n=== {d} (всего ttf: {len(files)}) ===")
            lines.append(f"Circe в этой папке: {circe or 'НЕТ'}")
    
    # Проверяем реальный путь установки LO
    result = subprocess.run(["find", "/usr", "-name", "*.ttf", "-path", "*libreoffice*"], 
                           capture_output=True, text=True)
    lo_ttfs = [l for l in result.stdout.splitlines() if "irce" in l]
    lines.append(f"\n=== Circe внутри LO (find) ===")
    lines.extend(lo_ttfs or ["НЕ НАЙДЕН"])
    
    return "\n".join(lines)
